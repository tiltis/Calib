"""기존 보고서 PPTX에 v9/v10/v11 버전 비교 슬라이드 추가"""
import shutil, os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy
from lxml import etree

SRC = r"C:\Users\tilti\AppData\Local\Temp\_report_tmp.pptx"
DST = r"C:\Users\tilti\OneDrive\samjung\output\report_orthorectification.pptx"

prs = Presentation(SRC)
W = prs.slide_width   # EMU
H = prs.slide_height

# ── 공통 스타일 헬퍼 ──
BG_COLOR   = RGBColor(0xFF, 0xFF, 0xFF)
TITLE_CLR  = RGBColor(0x1A, 0x1A, 0x2E)
SUB_CLR    = RGBColor(0x33, 0x33, 0x33)
ACCENT     = RGBColor(0x00, 0x70, 0xC0)
LIGHT_GRAY = RGBColor(0x88, 0x88, 0x88)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TBL_HEAD   = RGBColor(0x00, 0x56, 0x8A)
TBL_ALT    = RGBColor(0xF2, 0xF7, 0xFB)


def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def add_title(slide, text, sub=""):
    txBox = slide.shapes.add_textbox(Cm(1.5), Cm(0.8), Cm(22), Cm(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TITLE_CLR
    if sub:
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(14)
        p2.font.color.rgb = LIGHT_GRAY


def add_textbox(slide, left, top, width, height, texts, font_size=11):
    """texts: list of (text, bold, color)"""
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, bold, color) in enumerate(texts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(2)


def make_table(slide, left, top, width, rows_data, col_widths_cm):
    """rows_data: list of lists, first row = header"""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Cm(left), Cm(top),
                                        Cm(width), Cm(1.0 * n_rows))
    tbl = tbl_shape.table

    for j, w in enumerate(col_widths_cm):
        tbl.columns[j].width = Cm(w)

    for i, row in enumerate(rows_data):
        for j, cell_text in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = str(cell_text)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(10)
                para.font.color.rgb = WHITE if i == 0 else SUB_CLR
                para.font.bold = (i == 0)
                para.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 헤더 배경
            if i == 0:
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = TBL_HEAD
            elif i % 2 == 0:
                cell_fill = cell.fill
                cell_fill.solid()
                cell_fill.fore_color.rgb = TBL_ALT

    return tbl_shape


# ── 슬라이드 삽입 위치: Slide 3 뒤 (index 3) ──
INSERT_IDX = 3  # 0-based → slide 4 위치에 삽입


def insert_blank_slide(prs, idx):
    """빈 슬라이드를 idx 위치에 삽입"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    # 슬라이드를 원하는 위치로 이동
    slide_list = prs.slides._sldIdLst
    slide_elem = slide_list[-1]  # 방금 추가된 것 (맨 뒤)
    # idx 위치로 이동
    if idx < len(slide_list) - 1:
        ref_elem = slide_list[idx]
        slide_list.remove(slide_elem)
        slide_list.insert(idx, slide_elem)
    return slide


# ============================================================
# Slide A: 버전 진화 개요
# ============================================================
slideA = insert_blank_slide(prs, INSERT_IDX)
add_bg(slideA)
add_title(slideA, "알고리즘 버전 진화", "v9 (piza2) → v10 (piza3) → v11 (piza4)")

# 타임라인 식 버전 박스
versions = [
    ("v9 (piza2.py)", "Velocity 기반\nY축 비선형 보정\ncv2.remap()",
     "기준선 (안정)", RGBColor(0x4C, 0xAF, 0x50)),
    ("v10 (piza3.py)", "Y + X 동시 보정\n방사형 왜곡 제거\nscale_x = mpp/base",
     "2026-04-06", RGBColor(0x21, 0x96, 0xF3)),
    ("v11 (piza4.py)", "2차 다항 회귀\n단조성 강제\ncumsum 수정",
     "현업 권장", RGBColor(0xFF, 0x57, 0x22)),
]

for i, (title, desc, badge, color) in enumerate(versions):
    x = 1.5 + i * 7.5
    # 박스
    box = slideA.shapes.add_shape(1, Cm(x), Cm(4), Cm(6.5), Cm(8))  # 1=rectangle
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = ""
    p2.space_after = Pt(6)

    p3 = tf.add_paragraph()
    p3.text = desc
    p3.font.size = Pt(12)
    p3.font.color.rgb = WHITE
    p3.alignment = PP_ALIGN.CENTER
    p3.space_after = Pt(12)

    p4 = tf.add_paragraph()
    p4.text = badge
    p4.font.size = Pt(11)
    p4.font.bold = True
    p4.font.color.rgb = RGBColor(0xFF, 0xFF, 0x88)
    p4.alignment = PP_ALIGN.CENTER

    # 화살표 (마지막 제외)
    if i < 2:
        arrow = slideA.shapes.add_textbox(Cm(x + 6.5), Cm(7), Cm(1), Cm(2))
        atf = arrow.text_frame
        ap = atf.paragraphs[0]
        ap.text = "→"
        ap.font.size = Pt(30)
        ap.font.bold = True
        ap.font.color.rgb = LIGHT_GRAY
        ap.alignment = PP_ALIGN.CENTER

# 하단 요약
add_textbox(slideA, 1.5, 13, 22, 3, [
    ("핵심: 각 버전은 이전 버전의 모든 기능을 포함하며, 추가 보정을 적용합니다.", False, SUB_CLR),
    ("v11(piza4)이 현업 적용 권장 버전입니다 — 미지 파라미터 없이 측정 데이터만으로 동작.", True, ACCENT),
], font_size=11)


# ============================================================
# Slide B: v9→v10 상세 비교
# ============================================================
slideB = insert_blank_slide(prs, INSERT_IDX + 1)
add_bg(slideB)
add_title(slideB, "v9 → v10:  XY 동시 보정", "build_remap_tables()의 가로(X) 방향 보정 추가")

make_table(slideB, 1.5, 3.5, 22, [
    ["항목", "v9 (piza2)", "v10 (piza3)"],
    ["보정 방향", "세로(Y)만", "세로(Y) + 가로(X)"],
    ["map_x 계산", "map_x = col\n(가로 좌표 그대로)", "map_x = cx + (col-cx) / scale_x\n(행별 수평 스케일 적용)"],
    ["방사형 왜곡", "미처리\n→ 파도가 부채꼴로 잔존", "제거\n→ 먼 거리 파도를 물리적 폭으로 복원"],
    ["수식", "—", "scale_x(r) = mpp(src_y[r]) / base_mpp"],
    ["결과 비교", "위→아래 방사형 수렴", "전체 영역 등간격 격자"],
], col_widths_cm=[4, 8.5, 9.5])

add_textbox(slideB, 1.5, 12.5, 22, 4, [
    ("원리: 상단(먼 곳)의 행은 scale_x > 1 → 소스의 중앙 부근을 넓게 펼쳐 파도 마루 물리적 폭 복원", False, SUB_CLR),
    ("하단(가까운 곳)은 scale_x ≈ 1 → 변화 없음  |  등방성(isotropic) 원근 가정", False, LIGHT_GRAY),
], font_size=10)


# ============================================================
# Slide C: v10→v11 상세 비교
# ============================================================
slideC = insert_blank_slide(prs, INSERT_IDX + 2)
add_bg(slideC)
add_title(slideC, "v10 → v11:  회귀 모델 및 안정성 개선", "5개 핵심 변경으로 정확도/안정성 향상")

make_table(slideC, 1.5, 3.5, 22, [
    ["#", "변경 사항", "v10", "v11", "효과"],
    ["1", "속도 회귀 모델", "1차 (선형)", "1차/2차 자동 선택\n(R² 기준)", "R² 향상\n(비선형 포착)"],
    ["2", "mpp 단조성", "미보장", "누적 max 강제", "접힘 아티팩트\n방지"],
    ["3", "cumsum 시작", "mpp[0] (>0)", "0부터 시작", "상단 평탄 구간\n제거"],
    ["4", "출력 높이", "무제한", "roi_h × 3 상한", "메모리 폭발\n방지"],
    ["5", "영상 코덱", "mp4v 고정", "H264 우선 시도", "파일 크기\n50~70% 감소"],
], col_widths_cm=[1.2, 4.5, 4.5, 5.5, 4.3])

# 2차 회귀 설명 박스
add_textbox(slideC, 1.5, 12, 22, 5, [
    ("2차 회귀 자동 선택 기준:", True, ACCENT),
    ("  • 2차 R² - 1차 R² ≥ 0.02  →  2차 채택 (원근 비선형성 포착)", False, SUB_CLR),
    ("  • 2차 R² - 1차 R² < 0.02  →  1차 유지 (과적합 방지)", False, SUB_CLR),
    ("  • 단조성 강제: mpp = max_accumulate(mpp[::-1])[::-1]  →  위→아래 비증가 보장", False, SUB_CLR),
], font_size=10)


# ============================================================
# 저장
# ============================================================
prs.save(DST)
print(f"저장 완료: {DST}")
print(f"총 슬라이드: {len(prs.slides)}장 (기존 12 + 추가 3)")
