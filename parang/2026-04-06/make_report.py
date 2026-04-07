"""
make_report.py — 등거리 투영 보고서 PPTX 자동 생성

흐름:
  1. output/*.json 읽기 → 캘리브레이션 결과 수집
  2. 보정 전/후 좌우 비교 영상 생성 (opencv)
  3. python-pptx로 슬라이드 구성 + 비교 영상 삽입
  4. output/report_orthorectification.pptx 저장

저장 구조:
  output/
    {dir}_calib.json         ← piza4.py(v11) 또는 piza5.py(v12) 실행 후 생성됨
    {dir}_ortho.mp4          ← piza4.py(v11) 또는 piza5.py(v12) 실행 후 생성됨
    clips/{dir}_compare.mp4  ← 좌우 비교 영상 (이 스크립트가 생성)
    report_orthorectification.pptx
"""

import os
import json
import glob
import cv2
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm

# ============================================================
# 경로 설정
# ============================================================
OUTPUT_DIR  = r"C:\Users\tilti\OneDrive\samjung\output"
VIDEO_BASE  = r"C:\Users\tilti\OneDrive\samjung\data"
CLIP_DIR    = os.path.join(OUTPUT_DIR, "clips")
CLIP_SEC    = 15       # 비교 클립 길이 (초)
SLIDE_W     = Inches(13.33)
SLIDE_H     = Inches(7.5)

# ============================================================
# 색상 / 스타일
# ============================================================
C_NAVY   = RGBColor(0x1E, 0x3A, 0x5F)   # 제목 네이비
C_BLUE   = RGBColor(0x2E, 0x75, 0xB6)   # 강조 파랑
C_GRAY   = RGBColor(0x40, 0x40, 0x40)   # 본문 회색
C_LGRAY  = RGBColor(0xD9, 0xD9, 0xD9)   # 구분선
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK  = RGBColor(0x00, 0x00, 0x00)
FONT     = "맑은 고딕"


# ============================================================
# 유틸 — 좌우 비교 영상 생성 (Original | Ortho)
# ============================================================
TARGET_H = 480    # 비교 패널 높이

def make_compare_clip(orig_path, ortho_path, out_path, roi, seconds=CLIP_SEC, force=False):
    """
    원본(ROI 크롭) 좌측 + ortho 우측 합본 비교 영상 생성.
    하나의 영상에서 보정 전/후를 동시에 볼 수 있음.
    """
    if not force and os.path.exists(out_path) and os.path.getsize(out_path) > 0:
        print(f"  [비교 스킵] 이미 있음: {os.path.basename(out_path)}")
        return True
    if not os.path.exists(orig_path):
        print(f"  [비교 실패] 원본 없음")
        return False
    if not os.path.exists(ortho_path) or os.path.getsize(ortho_path) == 0:
        print(f"  [비교 실패] ortho 없음/손상")
        return False

    x1, y1, x2, y2 = roi
    roi_w, roi_h = x2 - x1, y2 - y1

    cap_orig  = cv2.VideoCapture(orig_path)
    cap_ortho = cv2.VideoCapture(ortho_path)
    fps = cap_orig.get(cv2.CAP_PROP_FPS) or 30.0
    total_cut = int(fps * seconds)

    ortho_w = int(cap_ortho.get(cv2.CAP_PROP_FRAME_WIDTH))
    ortho_h = int(cap_ortho.get(cv2.CAP_PROP_FRAME_HEIGHT))
    if ortho_w == 0 or ortho_h == 0:
        cap_orig.release(); cap_ortho.release()
        print(f"  [비교 실패] ortho 읽기 불가")
        return False

    # 패널 크기 계산 (둘 다 TARGET_H에 맞춤)
    orig_scale = TARGET_H / roi_h
    panel_w_L = int(roi_w * orig_scale)
    ortho_scale = TARGET_H / ortho_h
    panel_w_R = int(ortho_w * ortho_scale)

    gap = 4
    canvas_w = panel_w_L + gap + panel_w_R
    canvas_h = TARGET_H + 36  # 하단 라벨

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    fourcc = cv2.VideoWriter_fourcc(*"mp4v")
    writer = cv2.VideoWriter(out_path, fourcc, fps, (canvas_w, canvas_h))

    written = 0
    while written < total_cut:
        ret1, f_orig = cap_orig.read()
        ret2, f_ortho = cap_ortho.read()
        if not ret1 or not ret2:
            break

        f_orig = cv2.resize(f_orig, (1024, 576))
        roi_crop = f_orig[y1:y2, x1:x2]
        panel_L = cv2.resize(roi_crop, (panel_w_L, TARGET_H))
        panel_R = cv2.resize(f_ortho, (panel_w_R, TARGET_H))

        canvas = np.zeros((canvas_h, canvas_w, 3), dtype=np.uint8)
        canvas[:TARGET_H, :panel_w_L] = panel_L
        canvas[:TARGET_H, panel_w_L:panel_w_L+gap] = 255
        canvas[:TARGET_H, panel_w_L+gap:] = panel_R

        cv2.putText(canvas, "Original", (10, canvas_h - 10),
                    cv2.FONT_HERSHEY_SIMPLEX, 0.55, (200, 200, 200), 1)
        cv2.putText(canvas, "Ortho (v10)", (panel_w_L + gap + 10, canvas_h - 10),
                    cv2.FONT_HERSHEY_SIMPLEX, 0.55, (100, 255, 100), 1)

        writer.write(canvas)
        written += 1

    cap_orig.release()
    cap_ortho.release()
    writer.release()

    if written > 0:
        print(f"  [비교 생성] {os.path.basename(out_path)}  ({written/fps:.1f}s, {canvas_w}x{canvas_h})")
        return True
    else:
        if os.path.exists(out_path):
            os.remove(out_path)
        return False


# ============================================================
# 유틸 — 썸네일 추출 (비교 영상에서)
# ============================================================
def extract_thumbnail(video_path, out_path, at_sec=5):
    """영상 at_sec 초 지점 프레임을 PNG로 저장"""
    if os.path.exists(out_path):
        return out_path
    cap = cv2.VideoCapture(video_path)
    fps = cap.get(cv2.CAP_PROP_FPS) or 30.0
    cap.set(cv2.CAP_PROP_POS_FRAMES, int(fps * at_sec))
    ret, frame = cap.read()
    cap.release()
    if ret:
        cv2.imwrite(out_path, frame)
        return out_path
    return None


# ============================================================
# 유틸 — pptx 텍스트 박스
# ============================================================
def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=C_GRAY,
                align=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name      = font_name
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color):
    """배경 직사각형 추가"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# ============================================================
# 슬라이드 1 — 표지
# ============================================================
def make_title_slide(prs, dir_names):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # 상단 네이비 바
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.4), C_NAVY)

    # 제목
    add_textbox(slide, "해상 영상 등거리 투영 알고리즘",
                Inches(0.5), Inches(0.2), Inches(12), Inches(0.9),
                font_size=32, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    # 부제목
    add_textbox(slide, "Wave Orthorectification  |  보정 전/후 비교 보고서",
                Inches(0.5), Inches(1.5), Inches(10), Inches(0.5),
                font_size=16, color=C_BLUE)

    # 처리 영상 목록
    names_str = "\n".join(f"  \u2022 {n}" for n in dir_names) if dir_names else "  (처리된 결과 없음)"
    add_textbox(slide, f"처리 영상:\n{names_str}",
                Inches(0.5), Inches(2.2), Inches(8), Inches(3.5),
                font_size=14, color=C_GRAY)

    # 날짜 / 출처
    add_textbox(slide, "삼성중공업 CCTV 기반 파랑 계측\n2026.04",
                Inches(9.5), Inches(6.5), Inches(3.5), Inches(0.8),
                font_size=11, color=C_GRAY, align=PP_ALIGN.RIGHT)

    # 하단 라인
    add_rect(slide, 0, Inches(7.2), SLIDE_W, Pt(3), C_NAVY)

    return slide


# ============================================================
# 슬라이드 2 — 알고리즘 개요
# ============================================================
def make_algo_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 제목 바
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.75), C_NAVY)
    add_textbox(slide, "알고리즘 개요  \u2014  Velocity-based Orthorectification (v10)",
                Inches(0.3), Inches(0.1), Inches(12), Inches(0.55),
                font_size=20, bold=True, color=C_WHITE)

    # 흐름도 (텍스트 박스 5단계)
    steps = [
        ("\u2460 ROI 선택",        "마우스 드래그로\n해수면 영역 지정"),
        ("\u2461 파봉 추적 (20s)", "CLAHE\u2192Sobel\u2192peak\n프레임간 이동 추적"),
        ("\u2462 속도 맵 구축",    "슬라이딩 윈도우\n(y, px/s) 선형 회귀"),
        ("\u2463 m/px 보정",       "enc_speed / vel(y)\n= m/pixel(y)"),
        ("\u2464 등거리 투영",     "cv2.remap()으로\nXY 동시 비선형 보정"),
    ]
    box_w = Inches(2.3)
    box_h = Inches(1.8)
    gap   = Inches(0.15)
    top   = Inches(1.3)

    for i, (title, body) in enumerate(steps):
        lft = Inches(0.3) + i * (box_w + gap)
        # 배경 박스
        add_rect(slide, lft, top, box_w, box_h, RGBColor(0xE8, 0xF0, 0xFB))
        # 단계 제목
        add_textbox(slide, title, lft + Inches(0.08), top + Inches(0.08),
                    box_w - Inches(0.16), Inches(0.5),
                    font_size=13, bold=True, color=C_NAVY)
        # 본문
        add_textbox(slide, body, lft + Inches(0.08), top + Inches(0.6),
                    box_w - Inches(0.16), Inches(1.1),
                    font_size=11, color=C_GRAY)
        # 화살표 (마지막 제외)
        if i < len(steps) - 1:
            add_textbox(slide, "\u2192",
                        lft + box_w, top + Inches(0.7),
                        gap + Inches(0.1), Inches(0.5),
                        font_size=18, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

    # v10 개선 내용
    note = (
        "v10 개선: \u2460 XY 동시 보정 (가로 방향 방사형 왜곡 제거)  "
        "\u2461 Weather_Info.xlsx 자동 로딩  "
        "\u2462 Swell/Wind dominance 기반 encounter speed 자동 선택"
    )
    add_textbox(slide, note,
                Inches(0.3), Inches(3.4), Inches(12.7), Inches(0.5),
                font_size=11, color=C_BLUE)

    # 수식
    formula = "m/pixel(y)  =  encounter_speed  /  velocity(y)      encounter_speed = c_wave + SOG \u00b7 cos(\u03b8)"
    add_textbox(slide, formula,
                Inches(0.5), Inches(4.1), Inches(12), Inches(0.6),
                font_size=13, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    # 물리 가정 요약
    assume = (
        "가정:  \u2460 해수면 평면 (파고 무시)  "
        "\u2461 카메라 정지 (롤링/피칭 없음)  "
        "\u2462 Pinhole 카메라 (렌즈 왜곡 없음)  "
        "\u2463 encounter speed 일정"
    )
    add_textbox(slide, assume,
                Inches(0.5), Inches(4.9), Inches(12), Inches(0.5),
                font_size=10, color=RGBColor(0x80, 0x80, 0x80))

    add_rect(slide, 0, Inches(7.2), SLIDE_W, Pt(3), C_NAVY)
    return slide


# ============================================================
# 슬라이드 3 — 결과 요약 테이블
# ============================================================
def make_summary_slide(prs, calib_list):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.75), C_NAVY)
    add_textbox(slide, "처리 결과 요약",
                Inches(0.3), Inches(0.1), Inches(12), Inches(0.55),
                font_size=20, bold=True, color=C_WHITE)

    if not calib_list:
        add_textbox(slide, "처리된 결과가 없습니다.",
                    Inches(1), Inches(2), Inches(11), Inches(2),
                    font_size=16, color=C_GRAY, align=PP_ALIGN.CENTER)
        add_rect(slide, 0, Inches(7.2), SLIDE_W, Pt(3), C_NAVY)
        return slide

    # 테이블 헤더
    headers = ["영상", "Method", "Enc(m/s)", "R\u00b2", "base m/px", "출력 크기"]
    col_w   = [Inches(3.8), Inches(1.6), Inches(1.4), Inches(1.0), Inches(1.5), Inches(1.8)]
    row_h   = Inches(0.45)
    tbl_top = Inches(0.9)
    lft     = Inches(0.3)

    # 헤더 행
    x = lft
    for h, cw in zip(headers, col_w):
        add_rect(slide, x, tbl_top, cw, row_h, C_NAVY)
        add_textbox(slide, h, x + Inches(0.05), tbl_top + Inches(0.05),
                    cw - Inches(0.1), row_h - Inches(0.05),
                    font_size=11, bold=True, color=C_WHITE)
        x += cw

    # 데이터 행
    for ri, c in enumerate(calib_list):
        y   = tbl_top + row_h * (ri + 1)
        bg  = RGBColor(0xF2, 0xF7, 0xFF) if ri % 2 == 0 else C_WHITE
        x   = lft
        enc   = f"{c.get('enc_speed', 0):.2f}"
        r2    = f"{c.get('r_sq', 0):.3f}" if c.get('r_sq') is not None else "\u2014"
        mpp   = f"{c.get('base_mpp', 0):.3f}"
        sz    = f"{c['out_size'][0]}\u00d7{c['out_size'][1]}" if c.get('out_size') else "\u2014"
        meth  = "remap" if c.get('use_remap') else "homography"
        vals  = [c.get('dir_name',''), meth, enc, r2, mpp, sz]

        for v, cw in zip(vals, col_w):
            add_rect(slide, x, y, cw, row_h, bg)
            add_textbox(slide, str(v), x + Inches(0.05), y + Inches(0.05),
                        cw - Inches(0.1), row_h - Inches(0.05),
                        font_size=10, color=C_GRAY)
            x += cw

    add_rect(slide, 0, Inches(7.2), SLIDE_W, Pt(3), C_NAVY)
    return slide


# ============================================================
# 슬라이드 4+ — 영상별 상세 (파라미터 + 비교 영상)
# ============================================================
def make_video_slide(prs, calib, compare_clip, thumb_path):
    """
    왼쪽: 캘리브 파라미터 요약
    오른쪽: 좌우 비교 영상 (Original | Ortho) 한 개 — 동시 재생
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    dir_name = calib.get('dir_name', '')
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.75), C_NAVY)
    add_textbox(slide, f"결과: {dir_name}",
                Inches(0.3), Inches(0.1), Inches(12), Inches(0.55),
                font_size=18, bold=True, color=C_WHITE)

    # ---- 왼쪽: 파라미터 요약 ----
    wx = calib.get('weather', {})
    r2   = calib.get('r_sq')
    mpp  = calib.get('base_mpp', 0)
    enc  = calib.get('enc_speed', 0)
    enc_lbl = calib.get('enc_label', '')
    meth = "remap (비선형)" if calib.get('use_remap') else "homography (fallback)"
    sz   = calib.get('out_size', [0, 0])

    params_text = (
        f"SOG:      {wx.get('sog_knots', 0):.1f} kts\n"
        f"Heading:  {wx.get('heading_deg', 0):.1f}\u00b0\n"
        f"\n"
        f"Swell  Hs={wx.get('swell_hs', 0):.2f}m  "
        f"Tp={wx.get('swell_tp', 0):.2f}s\n"
        f"Wind   Hs={wx.get('wind_hs', 0):.2f}m  "
        f"Tp={wx.get('wind_tp', 0):.2f}s\n"
        f"\n"
        f"Enc speed:  {enc:.2f} m/s\n"
        f"              ({enc_lbl})\n"
        f"\n"
        f"R\u00b2:          {f'{r2:.3f}' if r2 is not None else '\u2014'}\n"
        f"base m/px:  {mpp:.4f}\n"
        f"방법:        {meth}\n"
        f"출력 크기:  {sz[0]}\u00d7{sz[1]}"
    )
    add_textbox(slide, params_text,
                Inches(0.3), Inches(0.9), Inches(4.2), Inches(5.8),
                font_size=12, color=C_GRAY)

    # ---- 오른쪽: 좌우 비교 영상 (한 개) ----
    vid_left = Inches(4.8)
    vid_top  = Inches(1.0)
    vid_w    = Inches(8.2)
    vid_h    = Inches(5.5)

    add_textbox(slide, "\u25b6 보정 전/후 동시 비교 (Original | Ortho v10)",
                vid_left, vid_top - Inches(0.3), vid_w, Inches(0.3),
                font_size=12, bold=True, color=C_NAVY)

    if compare_clip and os.path.exists(compare_clip) and os.path.getsize(compare_clip) > 0:
        try:
            slide.shapes.add_movie(
                compare_clip, vid_left, vid_top, vid_w, vid_h,
                mime_type='video/mp4'
            )
        except Exception as e:
            print(f"  [영상 삽입 실패] {e}")
            if thumb_path and os.path.exists(thumb_path):
                slide.shapes.add_picture(thumb_path, vid_left, vid_top, vid_w, vid_h)
    elif thumb_path and os.path.exists(thumb_path):
        slide.shapes.add_picture(thumb_path, vid_left, vid_top, vid_w, vid_h)
    else:
        add_rect(slide, vid_left, vid_top, vid_w, vid_h, RGBColor(0xCC, 0xCC, 0xCC))
        add_textbox(slide, "영상 없음",
                    vid_left, vid_top + vid_h // 2 - Inches(0.3), vid_w, Inches(0.6),
                    font_size=14, color=C_GRAY, align=PP_ALIGN.CENTER)

    add_rect(slide, 0, Inches(7.2), SLIDE_W, Pt(3), C_NAVY)
    return slide


# ============================================================
# 메인
# ============================================================
def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    os.makedirs(CLIP_DIR,   exist_ok=True)

    # ---- calib JSON 수집 ----
    json_files  = sorted(glob.glob(os.path.join(OUTPUT_DIR, "*_calib.json")))
    calib_list  = []
    for jf in json_files:
        with open(jf, encoding="utf-8") as f:
            calib_list.append(json.load(f))
    print(f"  calib JSON {len(calib_list)}개 발견")

    dir_names = [c.get('dir_name', '') for c in calib_list]

    # ---- 비교 영상 생성 ----
    clip_data = []   # (dir_name, compare_clip, thumb_path)
    for c in calib_list:
        dir_name   = c.get('dir_name', '')
        video_file = c.get('video', '')
        roi        = c.get('roi', [33, 6, 1019, 460])
        orig_src   = os.path.join(VIDEO_BASE, dir_name, video_file)
        ortho_src  = os.path.join(OUTPUT_DIR, f"{dir_name}_ortho.mp4")
        compare_path = os.path.join(CLIP_DIR, f"{dir_name}_compare.mp4")
        thumb_path   = os.path.join(CLIP_DIR, f"{dir_name}_compare_thumb.png")

        print(f"\n  [{dir_name}]")
        ok = make_compare_clip(orig_src, ortho_src, compare_path, roi, CLIP_SEC, force=True)

        # 썸네일 (비교 영상에서 추출)
        if ok and os.path.exists(compare_path):
            extract_thumbnail(compare_path, thumb_path)

        clip_data.append((dir_name, compare_path if ok else None, thumb_path))

    # ---- PPTX 생성 ----
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("\n  PPTX 슬라이드 생성 중...")

    make_title_slide(prs, dir_names)
    print("  슬라이드 1: 표지")

    make_algo_slide(prs)
    print("  슬라이드 2: 알고리즘 개요")

    make_summary_slide(prs, calib_list)
    print("  슬라이드 3: 결과 요약 테이블")

    for c, (dir_name, compare_clip, thumb_path) in zip(calib_list, clip_data):
        make_video_slide(prs, c, compare_clip, thumb_path)
        print(f"  슬라이드 +: {dir_name}")

    if not calib_list:
        print("  (calib 결과 없음 - 표지/개요/빈요약 3장만 생성)")

    out_path = os.path.join(OUTPUT_DIR, "report_orthorectification.pptx")
    prs.save(out_path)
    print(f"\n  [완료] {out_path}")
    print(f"  슬라이드 수: {len(prs.slides)}")


if __name__ == '__main__':
    main()
